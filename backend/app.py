import io
import os
import re
import zipfile
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Tuple

import cv2
import numpy as np
import pytesseract
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from PIL import Image
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

app = FastAPI(title="PPT Easy Backend v5")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition"],
)

@app.get("/")
def root():
    return {"ok": True, "service": "ppt-easy-backend-v5", "message": "Open /health or POST /api/convert"}

@app.get("/health")
def health():
    return {"ok": True, "service": "ppt-easy-backend-v5"}

@app.post("/api/convert")
async def convert(
    file: UploadFile = File(...),
    output: str = Form("excel"),
    lang: str = Form("eng+hin"),
    excel_mode: str = Form("table"),
    ppt_mode: str = Form("editable_text"),
    include_image: str = Form("false"),
):
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(400, "Only .pptx files supported")
    safe_base = clean_name(Path(file.filename).stem)
    tmpdir = tempfile.mkdtemp(prefix="ppt_easy_")
    ppt_path = Path(tmpdir) / file.filename
    ppt_path.write_bytes(await file.read())

    try:
        slides = extract_slide_images(ppt_path)
        if not slides:
            raise HTTPException(400, "PPTX me embedded slide images nahi mili.")

        results = []
        for s in slides:
            img = s["image"]
            full_ocr = ocr_words(img, lang)
            rows = detect_table_cells(img, lang) if excel_mode == "table" else []
            if not rows or len(rows) < 2:
                rows = layout_grid_rows(full_ocr["words"], img.width, col_count=18)
            title_lines = top_text_lines(full_ocr["lines"], img.height)
            results.append({
                "slide_no": s["slide_no"],
                "image": img,
                "rows": rows,
                "words": full_ocr["words"],
                "lines": full_ocr["lines"],
                "title": title_lines[0] if title_lines else f"Slide {s['slide_no']}",
            })

        include_bg = include_image.lower() == "true"
        made = []
        if output in ("excel", "both"):
            xlsx_path = Path(tmpdir) / f"{safe_base}_editable_table.xlsx"
            create_excel(results, xlsx_path)
            made.append(xlsx_path)
        if output in ("ppt", "both"):
            ppt_out = Path(tmpdir) / f"{safe_base}_editable.pptx"
            create_ppt(results, ppt_out, ppt_mode, include_bg)
            made.append(ppt_out)

        if output == "both":
            zip_path = Path(tmpdir) / f"{safe_base}_converted_outputs.zip"
            with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as z:
                for p in made:
                    z.write(p, p.name)
            return FileResponse(zip_path, filename=zip_path.name, media_type="application/zip")
        elif output == "ppt":
            return FileResponse(made[0], filename=made[0].name, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        else:
            return FileResponse(made[0], filename=made[0].name, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Conversion error: {e}")


def clean_name(s: str) -> str:
    s = re.sub(r"[^A-Za-z0-9_\-]+", "_", s or "converted").strip("_")
    return s[:80] or "converted"


def extract_slide_images(ppt_path: Path) -> List[Dict[str, Any]]:
    """Extract largest image from each slide. Works best for image-only PPT."""
    out = []
    with zipfile.ZipFile(ppt_path, "r") as z:
        names = set(z.namelist())
        slide_paths = sorted([n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n)], key=lambda x: int(re.search(r"slide(\d+)", x).group(1)))
        for slide_path in slide_paths:
            slide_no = int(re.search(r"slide(\d+)", slide_path).group(1))
            rel_path = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
            if rel_path not in names:
                continue
            slide_xml = z.read(slide_path).decode("utf-8", "ignore")
            rel_xml = z.read(rel_path).decode("utf-8", "ignore")
            rels = dict(re.findall(r'<Relationship[^>]*Id="([^"]+)"[^>]*Target="([^"]+)"', rel_xml))
            pics = []
            for m in re.finditer(r'<a:blip[^>]*(?:r:embed|r:link)="([^"]+)"[^>]*>', slide_xml):
                rid = m.group(1)
                around = slide_xml[max(0, m.start()-2500): m.end()+2500]
                ext = re.search(r'<a:ext[^>]*cx="(\d+)"[^>]*cy="(\d+)"', around)
                area = int(ext.group(1))*int(ext.group(2)) if ext else 0
                target = rels.get(rid, "")
                media_path = normalize_path("ppt/slides", target)
                if media_path in names and re.search(r"\.(png|jpe?g)$", media_path, re.I):
                    pics.append((area, media_path))
            if not pics:
                continue
            pics.sort(reverse=True)
            media_path = pics[0][1]
            img_bytes = z.read(media_path)
            img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
            out.append({"slide_no": slide_no, "image": img, "media_path": media_path})
    if not out:
        with zipfile.ZipFile(ppt_path, "r") as z:
            media = sorted([n for n in z.namelist() if re.match(r"ppt/media/.*\.(png|jpe?g)$", n, re.I)])
            for i, p in enumerate(media, start=1):
                img = Image.open(io.BytesIO(z.read(p))).convert("RGB")
                out.append({"slide_no": i, "image": img, "media_path": p})
    return out


def normalize_path(base: str, target: str) -> str:
    if not target:
        return ""
    if target.startswith("/"):
        return target[1:]
    parts = (base + "/" + target).split("/")
    stack = []
    for p in parts:
        if not p or p == ".":
            continue
        if p == "..":
            if stack:
                stack.pop()
        else:
            stack.append(p)
    return "/".join(stack)


def pil_to_cv(img: Image.Image) -> np.ndarray:
    return cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)


def ocr_words(img: Image.Image, lang: str) -> Dict[str, Any]:
    data = pytesseract.image_to_data(img, lang=lang, config="--oem 3 --psm 6", output_type=pytesseract.Output.DICT)
    words = []
    n = len(data.get("text", []))
    for i in range(n):
        text = str(data["text"][i] or "").strip()
        if not text:
            continue
        try:
            conf = float(data.get("conf", [0])[i])
        except Exception:
            conf = 0
        if conf < 15:
            continue
        x, y, w, h = int(data["left"][i]), int(data["top"][i]), int(data["width"][i]), int(data["height"][i])
        words.append({"text": text, "x0": x, "y0": y, "x1": x+w, "y1": y+h, "cx": x+w/2, "cy": y+h/2, "conf": conf})
    return {"words": words, "lines": group_words_into_lines(words)}


def group_words_into_lines(words: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    words = sorted(words, key=lambda w: (w["cy"], w["x0"]))
    if not words:
        return []
    hs = [max(8, w["y1"]-w["y0"]) for w in words]
    tol = max(8, float(np.median(hs))*0.75)
    groups = []
    for w in words:
        best = None; best_d = 10**9
        for g in groups:
            d = abs(g["cy"] - w["cy"])
            if d < tol and d < best_d:
                best = g; best_d = d
        if best is None:
            groups.append({"cy": w["cy"], "words": [w]})
        else:
            best["words"].append(w)
            best["cy"] = sum(x["cy"] for x in best["words"]) / len(best["words"])
    lines = []
    for g in groups:
        ws = sorted(g["words"], key=lambda w: w["x0"])
        text = " ".join(w["text"] for w in ws).strip()
        lines.append({
            "text": text,
            "words": ws,
            "x0": min(w["x0"] for w in ws), "y0": min(w["y0"] for w in ws),
            "x1": max(w["x1"] for w in ws), "y1": max(w["y1"] for w in ws),
            "cy": sum(w["cy"] for w in ws)/len(ws),
        })
    return sorted(lines, key=lambda l: (l["y0"], l["x0"]))


def top_text_lines(lines: List[Dict[str, Any]], h: int) -> List[str]:
    return [l["text"] for l in lines if l["y0"] < h*0.22 and len(l["text"]) > 3][:4]


def layout_grid_rows(words: List[Dict[str, Any]], img_w: int, col_count: int = 18) -> List[List[str]]:
    lines = group_words_into_lines(words)
    rows = []
    for line in lines:
        row = [""] * col_count
        for w in line["words"]:
            c = max(0, min(col_count-1, int(w["cx"] / max(1, img_w) * col_count)))
            row[c] = (row[c] + " " + w["text"]).strip()
        rows.append(row)
    return rows


def detect_table_cells(img: Image.Image, lang: str) -> List[List[str]]:
    """Find visible grid lines and OCR each cell. Fallback returns []."""
    cv = pil_to_cv(img)
    gray = cv2.cvtColor(cv, cv2.COLOR_BGR2GRAY)
    # upscale small images for better line detection
    scale = 2 if max(gray.shape) < 1800 else 1
    if scale > 1:
        gray = cv2.resize(gray, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)
    h, w = gray.shape
    # threshold: dark/colored line/text on light bg
    blur = cv2.GaussianBlur(gray, (3,3), 0)
    bw = cv2.adaptiveThreshold(blur, 255, cv2.ADAPTIVE_THRESH_MEAN_C, cv2.THRESH_BINARY_INV, 25, 10)
    # line kernels
    horiz_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (max(20, w//40), 1))
    vert_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, max(15, h//45)))
    horiz = cv2.morphologyEx(bw, cv2.MORPH_OPEN, horiz_kernel, iterations=1)
    vert = cv2.morphologyEx(bw, cv2.MORPH_OPEN, vert_kernel, iterations=1)
    grid = cv2.add(horiz, vert)
    contours, _ = cv2.findContours(grid, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    boxes = [cv2.boundingRect(c) for c in contours]
    boxes = [(x,y,bw_,bh_) for x,y,bw_,bh_ in boxes if bw_ > w*0.12 and bh_ > h*0.12]
    if boxes:
        # choose largest table-like area, but avoid whole slide border
        boxes.sort(key=lambda b: b[2]*b[3], reverse=True)
        x,y,tw,th = boxes[0]
    else:
        # fallback: detect densest lower area as table by projection
        x, y, tw, th = int(w*0.03), int(h*0.25), int(w*0.94), int(h*0.68)

    crop_h = horiz[y:y+th, x:x+tw]
    crop_v = vert[y:y+th, x:x+tw]
    # projections find line centers
    y_proj = crop_h.sum(axis=1) / 255
    x_proj = crop_v.sum(axis=0) / 255
    y_lines = merge_positions(np.where(y_proj > max(25, tw*0.22))[0].tolist(), gap=4)
    x_lines = merge_positions(np.where(x_proj > max(15, th*0.16))[0].tolist(), gap=4)
    # convert back to original scale coordinates
    y_lines = [int((p + y)/scale) for p in y_lines]
    x_lines = [int((p + x)/scale) for p in x_lines]
    # clean / include edges
    y_lines = clean_lines(y_lines, min_gap=8)
    x_lines = clean_lines(x_lines, min_gap=10)
    if len(y_lines) < 3 or len(x_lines) < 3:
        return []
    if len(x_lines) > 35:
        x_lines = reduce_lines(x_lines, 30)
    if len(y_lines) > 60:
        y_lines = reduce_lines(y_lines, 55)

    rows = []
    for r in range(len(y_lines)-1):
        row = []
        y0, y1 = y_lines[r], y_lines[r+1]
        if y1-y0 < 8:
            continue
        for c in range(len(x_lines)-1):
            x0, x1 = x_lines[c], x_lines[c+1]
            if x1-x0 < 8:
                row.append(""); continue
            pad = 2
            cell = img.crop((max(0,x0+pad), max(0,y0+pad), min(img.width,x1-pad), min(img.height,y1-pad)))
            txt = pytesseract.image_to_string(cell, lang=lang, config="--oem 3 --psm 6").strip()
            txt = re.sub(r"\s+", " ", txt)
            row.append(txt)
        # skip empty rows
        if any(cell for cell in row):
            rows.append(row)
    # remove too-empty columns
    rows = remove_empty_columns(rows)
    return rows


def merge_positions(pos: List[int], gap: int = 4) -> List[int]:
    if not pos:
        return []
    groups = [[pos[0]]]
    for p in pos[1:]:
        if p - groups[-1][-1] <= gap:
            groups[-1].append(p)
        else:
            groups.append([p])
    return [int(np.mean(g)) for g in groups]


def clean_lines(lines: List[int], min_gap: int) -> List[int]:
    lines = sorted(set(int(x) for x in lines))
    out = []
    for p in lines:
        if not out or p - out[-1] >= min_gap:
            out.append(p)
        else:
            out[-1] = int((out[-1] + p) / 2)
    return out


def reduce_lines(lines: List[int], max_count: int) -> List[int]:
    if len(lines) <= max_count:
        return lines
    idx = np.linspace(0, len(lines)-1, max_count).round().astype(int)
    return [lines[i] for i in idx]


def remove_empty_columns(rows: List[List[str]]) -> List[List[str]]:
    if not rows:
        return rows
    maxc = max(len(r) for r in rows)
    norm = [r + [""]*(maxc-len(r)) for r in rows]
    keep = []
    for c in range(maxc):
        nonempty = sum(1 for r in norm if str(r[c]).strip())
        if nonempty >= max(1, len(norm)//12):
            keep.append(c)
    if not keep:
        return rows
    return [[r[c] for c in keep] for r in norm]


def create_excel(results: List[Dict[str, Any]], out_path: Path):
    wb = Workbook()
    wb.remove(wb.active)
    blue = PatternFill("solid", fgColor="1F4E78")
    light = PatternFill("solid", fgColor="D9EAF7")
    white_font = Font(color="FFFFFF", bold=True)
    bold = Font(bold=True)
    side = Side(style="thin", color="BFBFBF")
    border = Border(left=side, right=side, top=side, bottom=side)
    for res in results:
        ws = wb.create_sheet(f"Slide {res['slide_no']}")
        rows = res["rows"] or [[l["text"]] for l in res["lines"]]
        maxc = max(1, max(len(r) for r in rows))
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=maxc)
        ws.cell(1,1).value = res.get("title") or f"Slide {res['slide_no']}"
        ws.cell(1,1).fill = blue
        ws.cell(1,1).font = white_font
        ws.cell(1,1).alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        start = 3
        for r_idx, row in enumerate(rows, start=start):
            for c_idx in range(1, maxc+1):
                v = row[c_idx-1] if c_idx-1 < len(row) else ""
                cell = ws.cell(r_idx, c_idx, v)
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                cell.border = border
                if r_idx == start or is_header_like(str(v)):
                    cell.fill = light
                    cell.font = bold
        for c in range(1, maxc+1):
            ws.column_dimensions[get_column_letter(c)].width = 15 if c < 4 else 11
        for r in range(1, ws.max_row+1):
            ws.row_dimensions[r].height = 24
        ws.freeze_panes = "A3"
    wb.save(out_path)


def is_header_like(s: str) -> bool:
    return bool(re.search(r"WORK|TYPE|GP|COMP|PHY|ONG|VERIF|TOT|SANC|BOOK|EXP|%|#|Rank|Score", s, re.I))


def create_ppt(results: List[Dict[str, Any]], out_path: Path, ppt_mode: str, include_image: bool):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for res in results:
        slide = prs.slides.add_slide(blank)
        if include_image:
            img_path = Path(tempfile.mkdtemp()) / f"slide_{res['slide_no']}.png"
            res["image"].save(img_path)
            slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        add_textbox(slide, "Converted Editable Output • SRDM ZP SATNA", 0.2, 0.08, 12.9, 0.28, 14, True, PP_ALIGN.CENTER, "0B3D73")
        if ppt_mode == "editable_table":
            add_table_slide(slide, res)
        else:
            add_layout_text(slide, res)
    prs.save(out_path)


def add_layout_text(slide, res: Dict[str, Any]):
    img = res["image"]
    W, H = 13.333, 7.5
    for l in res["lines"][:120]:
        x = max(0.05, min(W-0.1, l["x0"] / img.width * W))
        y = max(0.42, min(H-0.2, l["y0"] / img.height * H))
        w = max(0.25, min(W-x-0.05, (l["x1"]-l["x0"]) / img.width * W + 0.08))
        h = max(0.12, min(0.45, (l["y1"]-l["y0"]) / img.height * H + 0.04))
        fs = max(4.5, min(12, h * 52))
        add_textbox(slide, l["text"], x, y, w, h, fs, is_header_like(l["text"]), PP_ALIGN.LEFT, "111111")


def add_table_slide(slide, res: Dict[str, Any]):
    rows = res["rows"][:32]
    if not rows:
        add_layout_text(slide, res); return
    maxc = max(1, max(len(r) for r in rows))
    rows = [r + [""]*(maxc-len(r)) for r in rows]
    table_shape = slide.shapes.add_table(len(rows), maxc, Inches(0.15), Inches(0.55), Inches(13.05), Inches(6.45))
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val or "")
            cell.text_frame.margin_left = Pt(1)
            cell.text_frame.margin_right = Pt(1)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Nirmala UI"
                p.font.size = Pt(4 if maxc > 18 else 6)
                p.alignment = PP_ALIGN.CENTER
                if r == 0:
                    p.font.bold = True


def add_textbox(slide, text, x, y, w, h, font_size, bold=False, align=PP_ALIGN.LEFT, color="111111"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text or "")
    run.font.name = "Nirmala UI"
    run.font.size = Pt(font_size)
    run.font.bold = bool(bold)
    run.font.color.rgb = hex_to_rgb(color)
    p.alignment = align


def hex_to_rgb(hexstr: str):
    from pptx.dml.color import RGBColor
    s = hexstr.strip().lstrip('#')
    return RGBColor(int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))
